# backend/main.py
from fastapi import FastAPI
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from openai import OpenAI
from dotenv import load_dotenv
import os
import numpy as np
import glob
from pptx import Presentation
from pypdf import PdfReader
from docx import Document
from threading import Thread

# 환경 변수 로드
load_dotenv()

app = FastAPI()

# CORS 설정
app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://chatbot-frontend-plum-five.vercel.app"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# OpenAI 클라이언트
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# -------------------------------
# 📂 1. 문서별 텍스트 추출 함수
# -------------------------------
def load_pdf_text(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        if page.extract_text():
            text += page.extract_text() + "\n"
    return text

def load_ppt_text(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_docx_text(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

# -------------------------------
# ✂️ 2. 텍스트 청크 분할
# -------------------------------
def split_text(text, chunk_size=300):
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

# -------------------------------
# 🔎 3. 임베딩 함수
# -------------------------------
def get_embedding(text):
    response = client.embeddings.create(
        model="text-embedding-3-small",
        input=text
    )
    return response.data[0].embedding

# -------------------------------
# 🗄️ 4. 문서 로드 & 벡터화
# -------------------------------
def load_all_documents(folder_path="documents"):
    chunks = []

    # PDF
    for file in glob.glob(os.path.join(folder_path, "*.pdf")):
        print(f"📄 Loading PDF: {file}")
        chunks.extend(split_text(load_pdf_text(file)))

    # PPTX
    for file in glob.glob(os.path.join(folder_path, "*.pptx")):
        print(f"📊 Loading PPTX: {file}")
        chunks.extend(split_text(load_ppt_text(file)))

    # DOCX
    for file in glob.glob(os.path.join(folder_path, "*.docx")):
        print(f"📝 Loading DOCX: {file}")
        chunks.extend(split_text(load_docx_text(file)))

    return chunks

# -------------------------------
# 서버 초기화 플래그 & 전역 변수
# -------------------------------
chunks = []
embeddings = np.array([])
is_ready = False

def initialize_embeddings():
    global chunks, embeddings, is_ready
    print("📂 문서 로딩 및 벡터화 시작...")
    chunks = load_all_documents("documents")
    embeddings = np.array([get_embedding(c) for c in chunks])
    is_ready = True
    print(f"✅ 문서 로드 완료! 총 {len(chunks)}개 청크 벡터화")

# 백그라운드 스레드로 초기화
Thread(target=initialize_embeddings).start()

# -------------------------------
# 🔍 5. 검색 함수 (코사인 유사도)
# -------------------------------
def search(query, top_k=3):
    query_emb = np.array(get_embedding(query))
    similarities = np.dot(embeddings, query_emb) / (
        np.linalg.norm(embeddings, axis=1) * np.linalg.norm(query_emb)
    )
    top_indices = np.argsort(similarities)[::-1][:top_k]
    return [chunks[i] for i in top_indices]

# -------------------------------
# 🤖 6. 챗봇 API
# -------------------------------
@app.post("/chat")
async def chat(user_input: dict):
    global is_ready
    if not is_ready:
        return {"reply": "⚠️ 서버가 문서를 벡터화하는 중입니다. 잠시 후 다시 시도해주세요.", "context": []}

    query = user_input.get("message", "")
    if not query.strip():
        return {"reply": "질문을 입력해주세요.", "context": []}

    # 문서 검색
    context = search(query)

    # GPT에게 문맥과 함께 전달
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "너는 여러 문서를 기반으로 대답하는 AI야."},
            {"role": "user", "content": f"문서 내용: {context}\n\n질문: {query}"}
        ]
    )

    answer = response.choices[0].message.content
    return {"reply": answer, "context": context}

# -------------------------------
# 루트 테스트
# -------------------------------
@app.get("/")
def root():
    status_msg = "✅ 서버 실행 중"
    if not is_ready:
        status_msg = "⚠️ 문서 벡터화 진행 중"
    return {"message": f"FastAPI RAG 서버 정상 실행 🚀 - {status_msg}"}
